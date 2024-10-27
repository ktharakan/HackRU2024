from http.client import FORBIDDEN
import requests
from PyPDF2 import PdfReader, PdfMerger
from canvasapi import Canvas
from fpdf import FPDF
import os
from docx import Document
import pptx as ppt
import speech_recognition as sr


# Canvas API configuration
API_URL = "https://canvas.instructure.com"
API_KEY = ""
canvas = Canvas(API_URL, API_KEY)


r = sr.Recognizer()


def download_pdf(file_id, download_path):
    """Download a PDF file from Canvas."""
    try:
        file = canvas.get_file(file_id)
        file.download(download_path)
        print(f"Downloaded file: {download_path}")
    except Exception as e:
        print(f"Error downloading file: {str(e)}")


def clean_text(text):
    """Clean text by replacing problematic characters."""
    replacements = {
        '●': '*',  # Replace bullet points with asterisks
        '•': '*',  # Another type of bullet point
        '"': '"',  # Smart quotes
        '"': '"',
        ''': "'",
        ''': "'",
        '…': '...',
        '—': '-',
        '–': '-'
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text
pdf_counter = 1 
def create_pdf(input_file, display_name, template_pdf=None):
    """Create a PDF from text and optionally merge with a template."""
    global pdf_counter  # Use global counter
    try:
        # Ensure pdf_folder exists
        os.makedirs('pdf_folder', exist_ok=True)

        # Read and clean input text
        with open(input_file, 'r', encoding='utf-8') as f:
            text = f.read()
        
        text = clean_text(text)

        # Create PDF with cleaned text
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font('Helvetica', size=12)
        
        # Split text into paragraphs and write
        paragraphs = text.split('\n')
        for paragraph in paragraphs:
            if paragraph.strip():
                try:
                    pdf.write(5, paragraph + '\n')
                except Exception as e:
                    cleaned_paragraph = ''.join(char if ord(char) < 128 else '_' for char in paragraph)
                    pdf.write(5, cleaned_paragraph + '\n')

        # Use sequential naming for output
        output_path = os.path.join('pdf_folder', f"pdf{pdf_counter}.pdf")
        pdf.output(output_path)
        print(f"Created PDF: {output_path}")

        # If template exists, merge PDFs
        if template_pdf and os.path.exists(template_pdf):
            merger = PdfMerger()
            with open(template_pdf, 'rb') as template:
                merger.append(template)
            with open(output_path, 'rb') as output:
                merger.append(output)
            
            merged_path = os.path.join('pdf_folder', f"pdf{pdf_counter}.pdf")
            merger.write(merged_path)
            merger.close()
            print(f"Created merged PDF: {merged_path}")
        
        pdf_counter += 1  # Increment counter after successful creation
            
    except Exception as e:
        print(f"Error creating PDF: {str(e)}")


# Extract text from a PDF file
def convert_pdf_to_text(i, name):
    os.makedirs('unfiltered_folder', exist_ok=True)

    try:

        download_path = os.path.join('unfiltered_folder', name)
        
        i.download(download_path)
        pdf = PdfReader(download_path)

        text = []
        for i in pdf.pages:
            text+=i.extract_text()
        return ''.join(text)
    
    except:
        return ""
    
    
def docxtodoctotext(i, name):

    os.makedirs('docs_folder', exist_ok=True)


    try:
   
        download_path = os.path.join('docs_folder', name)
        i.download(download_path)
  
        document = Document(download_path)
        text = []
        for paragraph in document.paragraphs:
            text.append(paragraph.text)
        return '\n'.join(text)
    except Exception as e:
        print(f"Error processing .docx file: {str(e)}")
        return ""


def pptxtoppttotext(i, name):
    os.makedirs('ppt_folder', exist_ok=True)
    try:

        download_path = os.path.join('ppt_folder', name)
        i.download(download_path)
        prs = ppt.Presentation(
            download_path
        )
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except:
        return ""

def wavtotext(i, name):
    os.makedirs('wav_folder', exist_ok=True)
    try:
        download_path = os.path.join('wav_folder', i.display_name)
        i.download(download_path)
        prs = sr.AudioFile(download_path)
        with prs as source:
            audio=r.record(source)
            v = r.recognize_google(audio)
            return (str(v))
    except:
        return ""

canvas = Canvas(API_URL, API_KEY)


print("\nListing accessible courses:")
try:
    courses = canvas.get_courses(enrollment_state='active')
    for course in list(courses):
        print(f"Course ID: {course.id}, Name: {course.name}")
except Exception as e:
    print(f"Error listing courses: {str(e)}")

for course in list(courses):
  
    
    try:
        course_files = list(course.get_files())
    except Exception as e:
        print(f"Skipping course {course.name} - Unauthorized access")
        continue
    # Iterate through each file in the course
    for i in (list(course.get_files())):

        if(str(i.display_name).endswith(".docx") or str(i.display_name).endswith(".doc")):
            text = docxtodoctotext(i, i.display_name)
            
            # Create a  new text file each time
            with open('text.txt', 'w', encoding='utf-8') as f:

                f.write(text)
            create_pdf("text.txt", i.display_name,template_pdf='template.pdf' )
        
        if(str(i.display_name).endswith(".pdf")):
            text = convert_pdf_to_text(i, i.display_name)
            # Create a  new text file each time
            with open('text.txt', 'w', encoding='utf-8') as f:
                f.write(text)
            create_pdf("text.txt", i.display_name, template_pdf='template.pdf')

        if(str(i.display_name).endswith(".pptx") or str(i.display_name).endswith(".ppt")):

            text = pptxtoppttotext(i, i.display_name)

            # Create a  new text file each time
            with open('text.txt', 'w', encoding='utf-8') as f:
                f.write(text)

            create_pdf("text.txt", i.display_name,template_pdf='template.pdf' )

        if(str(i.display_name).endswith(".wav")):
            text = wavtotext(i, i.display_name)
            # Create a  new text file each time
            with open('text.txt', 'w', encoding='utf-8') as f:
                f.write(text)
            create_pdf("text.txt", i.display_name, template_pdf='template.pdf')



